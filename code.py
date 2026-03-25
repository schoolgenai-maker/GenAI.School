from pptx import Presentation
from pptx.util import Pt

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Define slide layouts
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(title_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "International School on Emerging Technologies of Generative AI"
    subtitle1.text = ("In Conjunction with the Workshop on AI for Healthcare @ENETCom\n"
                      "May 4th, 5th, & 6th (Revision Week)\n"
                      "ENETCom, University of Sfax, Tunisia")

    # --- Slide 2: General Information ---
    slide2 = prs.slides.add_slide(content_layout)
    title2 = slide2.shapes.title
    title2.text = "General Information & Audience"
    
    body2 = slide2.shapes.placeholders[1].text_frame
    body2.text = "Audience: PhD Students and Teachers"
    p = body2.add_paragraph()
    p.text = "Format: 5-hour sessions (9:00 AM – 2:00 PM)"
    p = body2.add_paragraph()
    p.text = "Philosophy: Fast-paced, bootcamp-style. Minimize slides, maximize coding and hands-on experimentation."
    p = body2.add_paragraph()
    p.text = "Contact: schoolgenai@gmail.com"

    # --- Slide 3: Speakers & Committee ---
    slide3 = prs.slides.add_slide(content_layout)
    title3 = slide3.shapes.title
    title3.text = "Speakers & Organizing Committee"
    
    body3 = slide3.shapes.placeholders[1].text_frame
    body3.text = "Suggested Speakers / Trainers:"
    p = body3.add_paragraph()
    p.text = "- Bassem Ben Hamed"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "- Nesrine Trabelsi"
    p.level = 1
    
    p = body3.add_paragraph()
    p.text = "Organizing Committee:"
    p = body3.add_paragraph()
    p.text = "Omar Cheikhrouhou, Olfa Gaddour, Mohamed Neji, Amira Echtioui, Ali Khalfallah, Hassen Ben Ayed, Wassim Jmal, Raida Hentati, Yassine Aydi, Tarek Abbes, Fathi Kallel."
    p.level = 1

    # --- Slide 4: Day 1 ---
    slide4 = prs.slides.add_slide(content_layout)
    title4 = slide4.shapes.title
    title4.text = "Day 1: Foundations & Mastering LLMs"
    
    body4 = slide4.shapes.placeholders[1].text_frame
    body4.text = "09:00 – 10:30 | Theory: From Transformers to LLMs"
    p = body4.add_paragraph()
    p.text = "Decoder-only architecture, tokens, context window. Pre-training vs post-training."
    p.level = 1
    
    p = body4.add_paragraph()
    p.text = "10:30 – 11:30 | Practice: Advanced Prompt Engineering"
    p = body4.add_paragraph()
    p.text = "Few-shot prompting, Chain of Thought (CoT), structured outputs (JSON)."
    p.level = 1
    
    p = body4.add_paragraph()
    p.text = "12:00 – 14:00 | Theory & Practice: Fine-tuning vs In-Context Learning"
    p = body4.add_paragraph()
    p.text = "PEFT/LoRA demo. Workshop: Quick fine-tuning of Mistral/Llama 3 (Colab/Unsloth)."
    p.level = 1

    # --- Slide 5: Day 2 ---
    slide5 = prs.slides.add_slide(content_layout)
    title5 = slide5.shapes.title
    title5.text = "Day 2: Retrieval-Augmented Generation (RAG)"
    
    body5 = slide5.shapes.placeholders[1].text_frame
    body5.text = "09:00 – 10:30 | Theory: The RAG Pipeline"
    p = body5.add_paragraph()
    p.text = "Embeddings, vector databases (Pinecone/Chroma). Mitigating hallucinations."
    p.level = 1
    
    p = body5.add_paragraph()
    p.text = "10:30 – 11:30 | Practice: My First RAG Pipeline"
    p = body5.add_paragraph()
    p.text = "PDF ingestion, chunking, and simple retrieval using LangChain."
    p.level = 1
    
    p = body5.add_paragraph()
    p.text = "12:00 – 14:00 | Theory & Practice: Advanced RAG"
    p = body5.add_paragraph()
    p.text = "Hybrid search, re-ranking, query transformation. Evaluation (RAGAS)."
    p.level = 1

    # --- Slide 6: Day 3 ---
    slide6 = prs.slides.add_slide(content_layout)
    title6 = slide6.shapes.title
    title6.text = "Day 3: The Era of AI Agents"
    
    body6 = slide6.shapes.placeholders[1].text_frame
    body6.text = "09:00 – 10:30 | Theory: Reasoning & Tools"
    p = body6.add_paragraph()
    p.text = "ReAct framework, Agent architecture (planning, memory, tools)."
    p.level = 1
    
    p = body6.add_paragraph()
    p.text = "10:30 – 11:30 | Practice: Tool Use & Function Calling"
    p = body6.add_paragraph()
    p.text = "Teaching an LLM to call Python functions (math/weather)."
    p.level = 1
    
    p = body6.add_paragraph()
    p.text = "12:00 – 14:00 | Final Application"
    p = body6.add_paragraph()
    p.text = "Using LangGraph or n8n. Workshop: Build a collaborative multi-agent team."
    p.level = 1

    # Save the presentation
    filename = "GenAI_Bootcamp_ENETCom.pptx"
    prs.save(filename)
    print(f"Success! PowerPoint presentation saved as '{filename}'")

if __name__ == "__main__":
    create_presentation()